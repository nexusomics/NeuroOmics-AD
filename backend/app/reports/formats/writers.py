"""CSV / Excel / Word / PowerPoint / PDF / HTML report writers."""
from __future__ import annotations

import io
import logging
from pathlib import Path

import pandas as pd

from app.reports.model import ReportData

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------
def write_csv(report: ReportData, out_dir: Path) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    paths = []
    for t in report.tables:
        p = out_dir / f"{_slug(t.name)}.csv"
        t.data.to_csv(p, index=False)
        paths.append(p)
    return paths


# ---------------------------------------------------------------------------
# Excel (multi-sheet workbook)
# ---------------------------------------------------------------------------
def write_xlsx(report: ReportData, out_dir: Path) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    p = out_dir / f"{_slug(report.title)}.xlsx"
    with pd.ExcelWriter(p, engine="openpyxl") as writer:
        for t in report.tables:
            t.data.to_excel(writer, sheet_name=_slug(t.name)[:31] or "sheet", index=False)
        if not report.tables:
            pd.DataFrame({"message": ["No tables generated"]}).to_excel(writer, index=False)
    return p


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------
def write_docx(report: ReportData, out_dir: Path, include_code: bool = False) -> Path:
    from docx import Document
    from docx.shared import Inches, Pt

    out_dir.mkdir(parents=True, exist_ok=True)
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(10.5)
    doc.add_heading(report.title, level=0)
    if report.subtitle:
        doc.add_paragraph(report.subtitle).italic = True
    doc.add_paragraph(f"Generated: {report.generated_at} · {report.tool_version}")

    def render_section(sec, level: int = 1) -> None:
        doc.add_heading(sec.title, level=level)
        for para in sec.paragraphs:
            doc.add_paragraph(para)
        for sub in sec.subsections:
            render_section(sub, level + 1)

    for sec in report.sections:
        render_section(sec)

    if report.tables:
        doc.add_heading("Tables", level=1)
        for t in report.tables:
            doc.add_heading(t.name, level=2)
            if t.caption:
                doc.add_paragraph(t.caption).italic = True
            _docx_table(doc, t.data)
    if report.figures:
        doc.add_heading("Figures", level=1)
        for fig in report.figures:
            p = Path(fig.path)
            if p.exists():
                doc.add_picture(str(p), width=Inches(6.0))
                doc.add_paragraph(fig.caption).italic = True
    if report.references:
        doc.add_heading("References", level=1)
        for ref in report.references:
            doc.add_paragraph(ref, style="List Number")
    out = out_dir / f"{_slug(report.title)}.docx"
    doc.save(out)
    return out


def _docx_table(doc, df: pd.DataFrame) -> None:
    from docx.shared import Pt

    table = doc.add_table(rows=df.shape[0] + 1, cols=df.shape[1])
    table.style = "Light Grid Accent 1"
    for j, col in enumerate(df.columns):
        table.rows[0].cells[j].text = str(col)
    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            v = df.iloc[i, j]
            table.rows[i + 1].cells[j].text = "" if pd.isna(v) else str(v)
    for row in table.rows:
        for cell in row.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(8)


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------
def write_pptx(report: ReportData, out_dir: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    slide_w = prs.slide_width
    # title slide
    title_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = report.title
    s.placeholders[1].text = report.subtitle or report.tool_version
    # content slides per section
    content_layout = prs.slide_layouts[1]
    for sec in report.sections:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = sec.title
        body = slide.placeholders[1].text_frame
        body.clear()
        for para in sec.paragraphs[:6]:
            p = body.add_paragraph()
            p.text = para[:220]
            p.font.size = Pt(14)
    # figures
    for fig in report.figures:
        p = Path(fig.path)
        if not p.exists():
            continue
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = fig.name
        try:
            slide.shapes.add_picture(str(p), Inches(1.2), Inches(1.6), width=Inches(7.6))
        except Exception:  # noqa: BLE001
            continue
    # tables
    for t in report.tables[:10]:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = t.name
        rows, cols = min(t.data.shape[0] + 1, 12), min(t.data.shape[1], 6)
        table = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.5), slide_w - Inches(0.8), Inches(3.5)).table
        for j, col in enumerate(t.data.columns[:6]):
            table.cell(0, j).text = str(col)
        for i in range(rows - 1):
            for j in range(cols):
                v = t.data.iloc[i, j] if i < t.data.shape[0] else ""
                table.cell(i + 1, j).text = "" if pd.isna(v) else str(v)[:30]
    out = out_dir / f"{_slug(report.title)}.pptx"
    prs.save(out)
    return out


# ---------------------------------------------------------------------------
# PDF (ReportLab)
# ---------------------------------------------------------------------------
def write_pdf(report: ReportData, out_dir: Path, dpi: int = 300) -> Path:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (
        Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
    )

    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / f"{_slug(report.title)}.pdf"
    styles = getSampleStyleSheet()
    body = ParagraphStyle("Body", parent=styles["BodyText"], fontSize=9.5, leading=13, spaceAfter=6)
    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontSize=15, spaceBefore=10, spaceAfter=6)
    caption = ParagraphStyle("Cap", parent=styles["Italic"], fontSize=8.5, textColor=colors.grey)

    story: list = [Paragraph(report.title, styles["Title"]), Spacer(1, 4)]
    story.append(Paragraph(f"Generated {report.generated_at} · {report.tool_version}", caption))
    story.append(Spacer(1, 8))
    for sec in report.sections:
        story.append(Paragraph(sec.title, h1))
        for para in sec.paragraphs:
            story.append(Paragraph(para, body))
    if report.figures:
        story.append(PageBreak())
        story.append(Paragraph("Figures", h1))
        for fig in report.figures:
            p = Path(fig.path)
            if p.exists():
                from PIL import Image as PILImage

                try:
                    w, h = PILImage.open(p).size
                    max_w = 6.3 * inch
                    ratio = max_w / w
                    story.append(Image(str(p), width=max_w, height=h * ratio))
                except Exception:  # noqa: BLE001
                    story.append(Image(str(p), width=6 * inch, height=4 * inch))
                story.append(Paragraph(f"<b>{fig.name}</b>. {fig.caption}", caption))
                story.append(Spacer(1, 10))
    for t in report.tables[:12]:
        story.append(Spacer(1, 6))
        story.append(Paragraph(t.name, h1))
        data = [t.data.columns.astype(str).tolist()]
        data += t.data.head(25).astype(str).values.tolist()
        tbl = Table(data, repeatRows=1)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f3b57")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTSIZE", (0, 0), (-1, -1), 7),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#eef2f6")]),
        ]))
        story.append(tbl)
    if report.references:
        story.append(PageBreak())
        story.append(Paragraph("References", h1))
        for ref in report.references:
            story.append(Paragraph(ref, body))
    doc = SimpleDocTemplate(str(out), pagesize=A4, rightMargin=0.8 * inch, leftMargin=0.8 * inch,
                            topMargin=0.7 * inch, bottomMargin=0.7 * inch)
    doc.build(story)
    return out


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------
def write_html(report: ReportData, out_dir: Path) -> Path:
    from jinja2 import Environment, FileSystemLoader, select_autoescape
    import base64

    out_dir.mkdir(parents=True, exist_ok=True)
    env = Environment(loader=FileSystemLoader(Path(__file__).parent / "templates"),
                      autoescape=select_autoescape(["html"]))
    figures = []
    for fig in report.figures:
        p = Path(fig.path)
        b64 = ""
        if p.exists():
            b64 = base64.b64encode(p.read_bytes()).decode()
        figures.append({"name": fig.name, "caption": fig.caption, "b64": b64, "ext": p.suffix.lstrip(".") if p.exists() else "png"})
    html = env.get_template("report.html.j2").render(
        report=report, figures=figures,
        tables=[{"name": t.name, "caption": t.caption, "html": t.data.head(50).to_html(classes="data-table", index=False)} for t in report.tables],
    )
    out = out_dir / f"{_slug(report.title)}.html"
    out.write_text(html, encoding="utf-8")
    return out


def _slug(text: str) -> str:
    keep = "".join(c if c.isalnum() or c in "-_ " else "_" for c in text)
    return "_".join(keep.strip().split()).lower()[:80] or "report"
